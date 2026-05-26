"""STELZ × Spot the Brand · pitch deck (22 mei 2026)

Cinematic JackandAI house style: black bg, off-white text, Spot Red accent.
Built specifically for tomorrow's STELZ meeting. Uses live data captured
21 May (Casa STELZ campaign discovery, lookalike algorithm proof, highlights).

Output: projects/spot-the-brand/STELZ-Pitch-Deck-2026-05-22.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_PATH = os.path.join(
    os.path.dirname(__file__), "..", "projects", "spot-the-brand",
    "STELZ-Pitch-Deck-2026-05-22.pptx",
)

# Brand colors — JackandAI cinematic
BLACK        = RGBColor(0x0A, 0x0A, 0x0A)
SURFACE      = RGBColor(0x14, 0x14, 0x14)
SURFACE_2    = RGBColor(0x1A, 0x1A, 0x1A)
BORDER       = RGBColor(0x2A, 0x2A, 0x2A)
OFFWHITE     = RGBColor(0xFF, 0xFD, 0xF3)
DIM          = RGBColor(0xBB, 0xBB, 0xBB)
MUTED        = RGBColor(0x77, 0x77, 0x77)
ACCENT       = RGBColor(0xFF, 0x13, 0x00)
GOOD         = RGBColor(0x4A, 0xDE, 0x80)
GOLD         = RGBColor(0xFB, 0xBF, 0x24)

FONT_DISPLAY = "Inter"      # Benzin substitute (system-safe)
FONT_MONO    = "JetBrains Mono"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL = 14


def set_bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=OFFWHITE, align=PP_ALIGN.LEFT, font=FONT_DISPLAY, italic=False, line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
        p.font.name = font
        p.alignment = align
        p.line_spacing = line_spacing
    return box


def add_rect(slide, left, top, width, height, color, *, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_outlined_rect(slide, left, top, width, height, line_color, line_pt=2, dash=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.background()
    shape.line.color.rgb = line_color
    shape.line.width = Pt(line_pt)
    if dash:
        shape.line.dash_style = 7
    return shape


def chrome(slide, n):
    add_text(slide, Inches(0.5), Inches(0.3), Inches(4), Inches(0.4),
             "Spot the Brand.", size=14, bold=True, color=OFFWHITE)
    add_rect(slide, Inches(2.06), Inches(0.55), Inches(0.08), Inches(0.08), ACCENT, rounded=True)
    add_text(slide, Inches(11.7), Inches(7.05), Inches(1.3), Inches(0.3),
             f"{n:02d} / {TOTAL:02d}", size=10, color=MUTED, align=PP_ALIGN.RIGHT, font=FONT_MONO)
    add_text(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
             "STELZ · 22 May 2026 · by JackandAI", size=10, color=MUTED, font=FONT_MONO)


def section_label(slide, text):
    add_text(slide, Inches(0.6), Inches(0.95), Inches(6), Inches(0.4),
             text, size=12, color=ACCENT, font=FONT_MONO, bold=True)


def crosshair(slide):
    add_rect(slide, Inches(0), Inches(3.75), SLIDE_W, Pt(0.5), BORDER)
    add_rect(slide, Inches(6.66), Inches(0), Pt(0.5), SLIDE_H, BORDER)


def stat_card(slide, x, y, w, h, label, value, sub=""):
    add_rect(slide, x, y, w, h, SURFACE)
    add_rect(slide, x, y, Inches(0.06), h, ACCENT)
    add_text(slide, x + Inches(0.3), y + Inches(0.25), w - Inches(0.4), Inches(0.4),
             label, size=10, color=MUTED, bold=True, font=FONT_MONO)
    add_text(slide, x + Inches(0.3), y + Inches(0.7), w - Inches(0.4), Inches(1.0),
             value, size=44, bold=True, color=OFFWHITE)
    if sub:
        add_text(slide, x + Inches(0.3), y + h - Inches(0.55), w - Inches(0.4), Inches(0.4),
                 sub, size=11, color=DIM)


def hit_card(slide, x, y, w, h, handle, followers, context, conf, platform="IG"):
    add_rect(slide, x, y, w, h, SURFACE)
    # Platform pill
    add_rect(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.45), Inches(0.25), SURFACE_2, rounded=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.21), Inches(0.45), Inches(0.25),
             platform, size=9, color=DIM, bold=True, align=PP_ALIGN.CENTER, font=FONT_MONO)
    # Confidence pill
    add_rect(slide, x + w - Inches(0.85), y + Inches(0.2), Inches(0.65), Inches(0.25), SURFACE_2, rounded=True)
    add_text(slide, x + w - Inches(0.85), y + Inches(0.21), Inches(0.65), Inches(0.25),
             conf, size=9, color=GOOD, bold=True, align=PP_ALIGN.CENTER, font=FONT_MONO)
    # Handle
    add_text(slide, x + Inches(0.2), y + Inches(0.65), w - Inches(0.4), Inches(0.4),
             handle, size=15, bold=True, color=OFFWHITE)
    add_text(slide, x + Inches(0.2), y + Inches(1.05), w - Inches(0.4), Inches(0.3),
             followers, size=10, color=MUTED, font=FONT_MONO)
    # Context
    add_text(slide, x + Inches(0.2), y + Inches(1.45), w - Inches(0.4), h - Inches(1.55),
             context, size=11, color=DIM, italic=True, line_spacing=1.25)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ════════════════════════════════════════════════════════
    # 01 · TITLE
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    crosshair(s)
    add_outlined_rect(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(2.5), ACCENT)
    add_text(s, Inches(1.0), Inches(2.25), Inches(6), Inches(0.3),
             "DETECTED · 1.00 · STELZ", size=11, color=ACCENT, font=FONT_MONO, bold=True)
    add_text(s, Inches(1.2), Inches(2.7), Inches(11), Inches(1.4),
             "STELZ × Spot the Brand.", size=72, bold=True, color=OFFWHITE)
    add_text(s, Inches(1.2), Inches(4.0), Inches(11), Inches(0.8),
             "Visual brand monitoring for the way people actually post.",
             size=22, color=OFFWHITE)
    add_text(s, Inches(1.2), Inches(5.3), Inches(11), Inches(0.5),
             "We see what Storyclash misses.",
             size=16, color=DIM, italic=True)
    add_text(s, Inches(0.5), Inches(0.3), Inches(4), Inches(0.4),
             "Spot the Brand.", size=14, bold=True, color=OFFWHITE)
    add_text(s, Inches(11.7), Inches(7.05), Inches(1.3), Inches(0.3),
             f"01 / {TOTAL:02d}", size=10, color=MUTED, align=PP_ALIGN.RIGHT, font=FONT_MONO)
    add_text(s, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
             "STELZ · 22 May 2026 · by JackandAI", size=10, color=MUTED, font=FONT_MONO)

    # ════════════════════════════════════════════════════════
    # 02 · THE BELIEF
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    chrome(s, 2)
    section_label(s, "WHAT WE BELIEVE")
    add_text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(2.6),
             "Brands live in the wild.\nNot in hashtags.",
             size=58, bold=True, color=OFFWHITE)
    add_text(s, Inches(0.6), Inches(4.5), Inches(12), Inches(0.6),
             "In someone's hand at een huisfeest. Op de bar in een story.",
             size=20, color=OFFWHITE)
    add_text(s, Inches(0.6), Inches(5.15), Inches(12), Inches(0.6),
             "Op een afstudeerbaret. Op een troon van blikjes. Op de muur in Casa STËLZ.",
             size=20, color=OFFWHITE)
    add_text(s, Inches(0.6), Inches(5.85), Inches(12), Inches(0.6),
             "Geen caption. Geen @drinkstelz. Gewoon STËLZ in beeld.",
             size=20, color=OFFWHITE)
    add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
             "Dat is waar jullie brand echt leeft. En dat is wat we zien.",
             size=14, color=MUTED, italic=True)

    # ════════════════════════════════════════════════════════
    # 03 · THE GAP
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    chrome(s, 3)
    section_label(s, "DE BLINDE VLEK")
    add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(1.0),
             "Wij zien beelden. Zij lezen tekst.",
             size=46, bold=True, color=OFFWHITE)

    add_rect(s, Inches(0.6), Inches(3.1), Inches(6.0), Inches(3.5), SURFACE)
    add_rect(s, Inches(0.6), Inches(3.1), Inches(0.06), Inches(3.5), MUTED)
    add_text(s, Inches(0.85), Inches(3.3), Inches(5.5), Inches(0.4),
             "STORYCLASH · MENTION · BRANDWATCH", size=10, color=MUTED, bold=True, font=FONT_MONO)
    add_text(s, Inches(0.85), Inches(3.75), Inches(5.5), Inches(0.6),
             "Captions, hashtags, @-mentions.", size=22, color=OFFWHITE)
    add_text(s, Inches(0.85), Inches(4.6), Inches(5.5), Inches(0.6),
             "Gebouwd toen mensen merken nog noemden.", size=14, color=DIM)
    add_text(s, Inches(0.85), Inches(5.55), Inches(5.5), Inches(1.0),
             "Vindt ~20% van brand mentions. De rest is onzichtbaar.",
             size=14, color=MUTED, italic=True)

    add_rect(s, Inches(6.85), Inches(3.1), Inches(6.0), Inches(3.5), SURFACE)
    add_rect(s, Inches(6.85), Inches(3.1), Inches(0.06), Inches(3.5), ACCENT)
    add_text(s, Inches(7.1), Inches(3.3), Inches(5.5), Inches(0.4),
             "SPOT THE BRAND", size=10, color=ACCENT, bold=True, font=FONT_MONO)
    add_text(s, Inches(7.1), Inches(3.75), Inches(5.5), Inches(0.6),
             "Visual AI op elk frame.", size=22, color=OFFWHITE)
    add_text(s, Inches(7.1), Inches(4.6), Inches(5.5), Inches(0.6),
             "Gemini Flash leest het blikje, de muur, de baret, het shirt.", size=14, color=DIM)
    add_text(s, Inches(7.1), Inches(5.55), Inches(5.5), Inches(1.0),
             "Vindt 100% van wat zichtbaar is. Inclusief de stille 80%.",
             size=14, color=OFFWHITE, italic=True)

    # ════════════════════════════════════════════════════════
    # 04 · LIVE STAND
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    chrome(s, 4)
    section_label(s, "LIVE STAND · 22 MEI 2026")
    add_text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.9),
             "Wat we nu al voor STELZ zien.",
             size=42, bold=True, color=OFFWHITE)

    stat_card(s, Inches(0.6), Inches(2.9), Inches(3.0), Inches(1.5), "CREATORS TRACKED", "3.047", "Instagram + TikTok")
    stat_card(s, Inches(3.8), Inches(2.9), Inches(3.0), Inches(1.5), "VISUAL HITS", "829", "in feed")
    stat_card(s, Inches(7.0), Inches(2.9), Inches(3.0), Inches(1.5), "CONFIRMED", "593", "door moderator")
    stat_card(s, Inches(10.2), Inches(2.9), Inches(2.7), Inches(1.5), "LAATSTE 24U", "+68", "nieuwe hits")

    stat_card(s, Inches(0.6), Inches(4.6), Inches(3.0), Inches(1.5), "PLATFORMS", "2", "IG · TikTok · Stories")
    stat_card(s, Inches(3.8), Inches(4.6), Inches(3.0), Inches(1.5), "PRODUCT LINES", "6", "alle smaken")
    stat_card(s, Inches(7.0), Inches(4.6), Inches(3.0), Inches(1.5), "SUBCULTURES", "1", "Casa STËLZ · live")
    stat_card(s, Inches(10.2), Inches(4.6), Inches(2.7), Inches(1.5), "FALSE POS RATE", "<2%", "mens-gekeurd")

    add_text(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.5),
             "Het systeem draait sinds februari. Wat hier staat is wat het automatisch heeft gevonden.",
             size=14, color=MUTED, italic=True)

    # ════════════════════════════════════════════════════════
    # 05 · CASA STËLZ STORY
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    chrome(s, 5)
    section_label(s, "CASE · CASA STËLZ")
    add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(1.0),
             "Gisteravond: één search.",
             size=46, bold=True, color=OFFWHITE)
    add_text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(0.5),
             "Vannacht: een hele campagne in jullie dashboard.",
             size=20, color=DIM, italic=True)

    # Timeline / process
    add_rect(s, Inches(0.6), Inches(3.3), Inches(12.1), Inches(2.6), SURFACE)
    steps = [
        ("21:00", "Search op TikTok: 'stelz house'"),
        ("21:05", "Casa STËLZ campagne ontdekt · niet in hashtag pool"),
        ("21:30", "#casastelz + #stelzhouse + #stelzibiza toegevoegd"),
        ("22:00", "Hashtag-scrape: 393 TikToks, 91 nieuwe creators"),
        ("22:30", "17 visual hits gedetecteerd · waaronder DJ met 1.8M likes"),
        ("23:00", "Lookalike-algoritme: 17 kandidaten in subculture"),
    ]
    for i, (t, line) in enumerate(steps):
        y = Inches(3.5 + i * 0.4)
        add_text(s, Inches(0.85), y, Inches(0.9), Inches(0.35), t, size=11, color=ACCENT, font=FONT_MONO, bold=True)
        add_rect(s, Inches(1.75), y + Inches(0.13), Inches(0.06), Inches(0.06), DIM)
        add_text(s, Inches(2.0), y, Inches(10), Inches(0.35), line, size=14, color=OFFWHITE)

    add_text(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.5),
             "Drie uur. Twee koffie. Géén handmatige creator-research.",
             size=15, color=DIM, italic=True)
    add_text(s, Inches(0.6), Inches(6.75), Inches(12), Inches(0.3),
             "Vraag aan STËLZ: welke campagnes hebben jullie nog meer?",
             size=13, color=ACCENT, font=FONT_MONO, bold=True)

    # ════════════════════════════════════════════════════════
    # 06 · GOLDEN VISUAL-ONLY HITS
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    chrome(s, 6)
    section_label(s, "VISUAL ONLY · WHAT OTHERS MISS")
    add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(1.0),
             "Geen caption. Geen hashtag. Pure organische zichtbaarheid.",
             size=28, bold=True, color=OFFWHITE)

    # 4 hit cards
    hit_card(s, Inches(0.6), Inches(3.0), Inches(3.0), Inches(3.7),
             "@niekroozen", "260.300 followers · 27.4K likes",
             '"Tour de house!" — STËLZ logo aan witte muur achter hem. Geen caption-mention.',
             "100%", "TT")
    hit_card(s, Inches(3.8), Inches(3.0), Inches(3.0), Inches(3.7),
             "@bavaria.bierkoerier", "3.537 followers",
             "Man zit op een TROON gebouwd van STËLZ Hard Seltzer blikjes. Onder gekleurd licht. Iconic.",
             "100%", "IG")
    hit_card(s, Inches(7.0), Inches(3.0), Inches(3.0), Inches(3.7),
             "@laurebaele", "739 followers · 400 likes",
             '"Leukste afsluit van het studentenleven #graduated" — STËLZ S-logo op haar afstudeerbaret.',
             "90%", "TT")
    hit_card(s, Inches(10.2), Inches(3.0), Inches(2.7), Inches(3.7),
             "@chloeannamarianne", "62 followers · stagiair",
             '"Lunch for @elle_nl" — STËLZ blikjes op de dining tafel. Gevonden via lookalike-algoritme.',
             "80%", "IG")

    # ════════════════════════════════════════════════════════
    # 07 · LOOKALIKE ALGORITME
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    chrome(s, 7)
    section_label(s, "DE OLIEVLEK")
    add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(1.0),
             "Geef ons jullie campagne. Wij vinden de scene.",
             size=38, bold=True, color=OFFWHITE)

    add_text(s, Inches(0.6), Inches(2.6), Inches(12.5), Inches(0.5),
             "Het lookalike-algoritme in vier stappen:",
             size=16, color=DIM)

    # 4 steps
    step_w = Inches(3.0)
    step_y = Inches(3.3)
    steps_lookalike = [
        ("01", "TAG", "Markeer jullie campagne-creators als subculture"),
        ("02", "MINE", "Systeem extract hun signature hashtags + co-occurrences"),
        ("03", "SCAN", "Apify-scrape onder die hashtags · nieuwe handles gevonden"),
        ("04", "SCORE", "Signal-score per kandidaat · drempel = kandidaat"),
    ]
    for i, (n, t, body) in enumerate(steps_lookalike):
        x = Inches(0.6 + i * 3.2)
        add_rect(s, x, step_y, step_w, Inches(2.3), SURFACE)
        add_rect(s, x, step_y, Inches(0.06), Inches(2.3), ACCENT)
        add_text(s, x + Inches(0.3), step_y + Inches(0.2), step_w - Inches(0.4), Inches(0.4),
                 n, size=10, color=ACCENT, font=FONT_MONO, bold=True)
        add_text(s, x + Inches(0.3), step_y + Inches(0.55), step_w - Inches(0.4), Inches(0.5),
                 t, size=22, bold=True, color=OFFWHITE)
        add_text(s, x + Inches(0.3), step_y + Inches(1.1), step_w - Inches(0.4), Inches(1.2),
                 body, size=12, color=DIM, line_spacing=1.25)

    # Proof badge
    add_rect(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(1.0), SURFACE_2)
    add_rect(s, Inches(0.6), Inches(5.9), Inches(0.06), Inches(1.0), GOOD)
    add_text(s, Inches(0.85), Inches(6.0), Inches(2.5), Inches(0.4),
             "BEWIJS · CASA STËLZ", size=10, color=GOOD, font=FONT_MONO, bold=True)
    add_text(s, Inches(0.85), Inches(6.35), Inches(12), Inches(0.5),
             "Binnen 1 uur na expansion: 3 van 17 lookalikes hadden al een visual hit.",
             size=16, color=OFFWHITE, bold=True)
    add_text(s, Inches(0.85), Inches(6.65), Inches(12), Inches(0.4),
             "@chloeannamarianne (signal 12) → blijkt stagiair bij Booij Agency. STËLZ in een Elle NL post.",
             size=12, color=DIM, italic=True)

    # ════════════════════════════════════════════════════════
    # 08 · STORIES & HIGHLIGHTS
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    chrome(s, 8)
    section_label(s, "BEYOND THE FEED")
    add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(1.0),
             "Stories. Highlights. Het brand op een baret.",
             size=38, bold=True, color=OFFWHITE)

    add_text(s, Inches(0.6), Inches(2.5), Inches(12), Inches(0.6),
             "Niet alleen feed posts. Wij scannen elk frame waar STELZ in beeld kan zijn.",
             size=16, color=DIM)

    # 3 columns
    cols = [
        ("HIGHLIGHTS", "440 frames", "Permanente saved stories. 21 STELZ hits in eerste run. Bavaria troon. Retail visibility. DJ booths.", GOLD),
        ("STORIES (24U)", "Cron-ready", "Live story-scraping via gekoppeld IG-account. Setup: 1 dag na deal. Verwachte uplift: 2-3× brand impressies.", ACCENT),
        ("TIKTOK", "194 hits", "Inclusief #casastelz, #stelzhouse, #stelzibiza. Cover-frame detection via Gemini Vision.", DIM),
    ]
    for i, (label, val, body, color) in enumerate(cols):
        x = Inches(0.6 + i * 4.25)
        add_rect(s, x, Inches(3.4), Inches(4.0), Inches(3.2), SURFACE)
        add_rect(s, x, Inches(3.4), Inches(0.06), Inches(3.2), color)
        add_text(s, x + Inches(0.3), Inches(3.6), Inches(3.5), Inches(0.4),
                 label, size=11, color=color, font=FONT_MONO, bold=True)
        add_text(s, x + Inches(0.3), Inches(4.0), Inches(3.5), Inches(0.8),
                 val, size=32, bold=True, color=OFFWHITE)
        add_text(s, x + Inches(0.3), Inches(5.0), Inches(3.5), Inches(1.7),
                 body, size=12, color=DIM, line_spacing=1.3)

    # ════════════════════════════════════════════════════════
    # 09 · HOW IT WORKS
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    chrome(s, 9)
    section_label(s, "HOW IT WORKS")
    add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(1.0),
             "Vier lagen. Eén dashboard.",
             size=42, bold=True, color=OFFWHITE)

    layers = [
        ("HARVEST", "Apify scrapers · IG hashtag, IG profile, IG stories, TikTok hashtag, TikTok profile"),
        ("DETECT", "Gemini 2.5 Flash op cover-image · 30-67 img/min · €0.001 per detection"),
        ("VERIFY", "Mens-in-the-loop moderator · sticky reject · trains de scanner"),
        ("SURFACE", "Live dashboard · filters · subcultures · CSV export · alerts via mail/Slack/WA"),
    ]
    for i, (t, body) in enumerate(layers):
        y = Inches(3.0 + i * 0.95)
        add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.8), SURFACE)
        add_rect(s, Inches(0.6), y, Inches(0.06), Inches(0.8), ACCENT)
        add_text(s, Inches(0.9), y + Inches(0.1), Inches(2.5), Inches(0.5),
                 t, size=18, bold=True, color=OFFWHITE)
        add_text(s, Inches(3.4), y + Inches(0.15), Inches(9.2), Inches(0.5),
                 body, size=13, color=DIM, line_spacing=1.2)

    # ════════════════════════════════════════════════════════
    # 10 · WHY NOW
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    chrome(s, 10)
    section_label(s, "WHY NOW")
    add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(1.0),
             "Drie shifts maken visual brand monitoring nu mogelijk.",
             size=30, bold=True, color=OFFWHITE)

    shifts = [
        ("01", "Mensen taggen niet meer", "Gen Z noemt brands niet in captions. 80% van brand visibility is visueel-only. Storyclash mist het."),
        ("02", "Vision AI is goedkoop", "Gemini Flash: €0.001 per image. Wat 5 jaar geleden niet kon, kost nu minder dan een hashtag-scrape."),
        ("03", "Stories zijn ephemeral", "24 uur visibility. Zonder daily scrape gemist. Onze cron pakt elke dag. Highlights blijven, stories ook gevangen."),
    ]
    for i, (n, t, body) in enumerate(shifts):
        y = Inches(3.0 + i * 1.25)
        add_text(s, Inches(0.6), y, Inches(1.0), Inches(0.6),
                 n, size=42, bold=True, color=ACCENT, font=FONT_MONO)
        add_text(s, Inches(1.7), y, Inches(11), Inches(0.5),
                 t, size=22, bold=True, color=OFFWHITE)
        add_text(s, Inches(1.7), y + Inches(0.55), Inches(11), Inches(0.6),
                 body, size=13, color=DIM, line_spacing=1.25)

    # ════════════════════════════════════════════════════════
    # 11 · PACKAGING
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    chrome(s, 11)
    section_label(s, "PACKAGING")
    add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(1.0),
             "Drie tiers. Eén platform.",
             size=42, bold=True, color=OFFWHITE)

    tiers = [
        ("PILOT", "€500/mnd", "3 maanden", [
            "Tot 500 creators",
            "1 platform (IG of TT)",
            "Weekly digest mail",
            "Mens-moderatie inbegrepen",
        ], DIM),
        ("GROWTH", "€1.500/mnd", "Most popular", [
            "Tot 3.000 creators",
            "IG + TikTok + Stories",
            "Daily alerts (mail/Slack/WA)",
            "Subculture expansion",
            "Lookalike algoritme",
        ], ACCENT),
        ("AGENCY", "€4.000/mnd", "Multi-brand", [
            "Onbeperkt creators",
            "Multi-brand dashboards",
            "API access · CSV exports",
            "Custom integraties",
            "Whitelabel optie",
        ], GOLD),
    ]
    for i, (name, price, sub, items, color) in enumerate(tiers):
        x = Inches(0.6 + i * 4.25)
        add_rect(s, x, Inches(2.7), Inches(4.0), Inches(4.0), SURFACE)
        add_rect(s, x, Inches(2.7), Inches(4.0), Inches(0.06), color)
        add_text(s, x + Inches(0.3), Inches(2.95), Inches(3.5), Inches(0.4),
                 name, size=12, color=color, font=FONT_MONO, bold=True)
        add_text(s, x + Inches(0.3), Inches(3.35), Inches(3.5), Inches(0.7),
                 price, size=30, bold=True, color=OFFWHITE)
        add_text(s, x + Inches(0.3), Inches(4.0), Inches(3.5), Inches(0.35),
                 sub, size=11, color=DIM, italic=True)
        for j, item in enumerate(items):
            add_text(s, x + Inches(0.3), Inches(4.5 + j * 0.4), Inches(3.5), Inches(0.35),
                     "· " + item, size=12, color=OFFWHITE)

    # ════════════════════════════════════════════════════════
    # 12 · THE 30-60-90
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    chrome(s, 12)
    section_label(s, "WAT JULLIE KRIJGEN")
    add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(1.0),
             "30 dagen. 60 dagen. 90 dagen.",
             size=42, bold=True, color=OFFWHITE)

    timeline = [
        ("DAG 1-30", "BASELINE",
         "3000+ creators in tracker · 1000+ visual hits · alle product-lines gemonitord · weekly cultuurrapport",
         "Voor jullie marketing & merk-team"),
        ("DAG 31-60", "ACTIVATIE",
         "Casa STËLZ + Lakedance + Solar als subcultures · lookalike-uitbreiding · creator outreach via DM-templates",
         "Voor jullie influencer programma"),
        ("DAG 61-90", "INTELLIGENCE",
         "Cultuurrapport per maand · sentiment scoring · POS shelf visibility · competitive benchmark",
         "Voor jullie commercial/sales team"),
    ]
    for i, (phase, t, body, who) in enumerate(timeline):
        y = Inches(3.0 + i * 1.2)
        add_rect(s, Inches(0.6), y, Inches(12.1), Inches(1.0), SURFACE)
        add_rect(s, Inches(0.6), y, Inches(0.06), Inches(1.0), ACCENT)
        add_text(s, Inches(0.85), y + Inches(0.1), Inches(2.0), Inches(0.4),
                 phase, size=11, color=ACCENT, font=FONT_MONO, bold=True)
        add_text(s, Inches(0.85), y + Inches(0.5), Inches(2.0), Inches(0.4),
                 t, size=18, bold=True, color=OFFWHITE)
        add_text(s, Inches(3.0), y + Inches(0.12), Inches(7.5), Inches(0.5),
                 body, size=12, color=OFFWHITE, line_spacing=1.3)
        add_text(s, Inches(3.0), y + Inches(0.62), Inches(7.5), Inches(0.4),
                 who, size=10, color=MUTED, italic=True, font=FONT_MONO)

    # ════════════════════════════════════════════════════════
    # 13 · LIVE NU
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    chrome(s, 13)
    section_label(s, "LIVE NU")
    add_text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(1.0),
             "Het draait al. Voor jullie.",
             size=46, bold=True, color=OFFWHITE)

    add_text(s, Inches(0.6), Inches(2.7), Inches(12), Inches(0.6),
             "spotyourbrand.com/?brand=stelz",
             size=24, color=ACCENT, font=FONT_MONO, bold=True)

    # Two cols: dashboard + moderator
    add_rect(s, Inches(0.6), Inches(3.6), Inches(6.0), Inches(2.8), SURFACE)
    add_rect(s, Inches(0.6), Inches(3.6), Inches(0.06), Inches(2.8), ACCENT)
    add_text(s, Inches(0.85), Inches(3.8), Inches(5.5), Inches(0.4),
             "DASHBOARD", size=12, color=ACCENT, font=FONT_MONO, bold=True)
    add_text(s, Inches(0.85), Inches(4.2), Inches(5.5), Inches(1.0),
             "829 hits live. Filters per product line, size, platform, format. Subcultures tab.",
             size=15, color=OFFWHITE, line_spacing=1.3)
    add_text(s, Inches(0.85), Inches(5.5), Inches(5.5), Inches(0.6),
             "CSV export. Daily delta-mail. Alle creator profielen klikbaar.",
             size=12, color=DIM)

    add_rect(s, Inches(6.85), Inches(3.6), Inches(6.0), Inches(2.8), SURFACE)
    add_rect(s, Inches(6.85), Inches(3.6), Inches(0.06), Inches(2.8), GOOD)
    add_text(s, Inches(7.1), Inches(3.8), Inches(5.5), Inches(0.4),
             "MODERATOR", size=12, color=GOOD, font=FONT_MONO, bold=True)
    add_text(s, Inches(7.1), Inches(4.2), Inches(5.5), Inches(1.0),
             "516 confirmed. 175 rejected. Sticky: reject = weg, ook bij rescan.",
             size=15, color=OFFWHITE, line_spacing=1.3)
    add_text(s, Inches(7.1), Inches(5.5), Inches(5.5), Inches(0.6),
             "Keyboard shortcuts. Trainen visie. False positive < 2%.",
             size=12, color=DIM)

    add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
             "Bekijk het tijdens deze meeting · live · met jullie ogen erbij.",
             size=14, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════
    # 14 · CLOSE / CTA
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    crosshair(s)
    add_text(s, Inches(0.5), Inches(0.3), Inches(4), Inches(0.4),
             "Spot the Brand.", size=14, bold=True, color=OFFWHITE)
    add_rect(s, Inches(2.06), Inches(0.55), Inches(0.08), Inches(0.08), ACCENT, rounded=True)
    add_text(s, Inches(11.7), Inches(7.05), Inches(1.3), Inches(0.3),
             f"14 / {TOTAL:02d}", size=10, color=MUTED, align=PP_ALIGN.RIGHT, font=FONT_MONO)
    add_text(s, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
             "STELZ · 22 May 2026 · by JackandAI", size=10, color=MUTED, font=FONT_MONO)

    add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(2.0),
             "Klein begin. Grote vlek.",
             size=72, bold=True, color=OFFWHITE)
    add_text(s, Inches(0.6), Inches(3.6), Inches(12), Inches(0.6),
             "Vandaag tekenen → morgen nieuwe campagne live in dashboard.",
             size=22, color=DIM)

    add_outlined_rect(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.6), ACCENT)
    add_text(s, Inches(0.85), Inches(5.2), Inches(11), Inches(0.4),
             "NEXT STEP", size=11, color=ACCENT, font=FONT_MONO, bold=True)
    add_text(s, Inches(0.85), Inches(5.55), Inches(11), Inches(0.6),
             "Pilot · €500/mnd · 3 maanden · binnen 48u live met jullie data.",
             size=22, bold=True, color=OFFWHITE)
    add_text(s, Inches(0.85), Inches(6.15), Inches(11), Inches(0.4),
             "Contract op één A4. Geen setup fee. Maand-tot-maand opzegbaar na pilot.",
             size=13, color=DIM)

    add_text(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.4),
             "Meinte Stinstra · meinte.stinstra@dorstenlesser.nl",
             size=12, color=MUTED, font=FONT_MONO, align=PP_ALIGN.CENTER)

    # Save
    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Saved: {OUTPUT_PATH}")


if __name__ == "__main__":
    build()
