"""Spot the Brand · pitch deck v0.2 — brand-led.

Re-ordered from v0.1 so the deck argues for the BRAND, not for STELZ.
Slides 2-4 are the manifesto thesis; slide 6 is proof where STELZ appears
as one example. STELZ is no longer the protagonist of the deck.

Visual identity: Spot the Brand · JackandAI house. Black bg, off-white
text, Spot Red (#FF1300) accent. Bounding-box + crosshair motifs.

Output: projects/spot-the-brand/Spot the Brand - Pitch Deck v0.pptx
(overwrites v0.1; that file is in git history if needed)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_PATH = os.path.join(
    os.path.dirname(__file__), "..", "projects", "spot-the-brand",
    "Spot the Brand - Pitch Deck v0.pptx",
)

# Brand colors
BLACK        = RGBColor(0x0A, 0x0A, 0x0A)
SURFACE      = RGBColor(0x14, 0x14, 0x14)
BORDER       = RGBColor(0x2A, 0x2A, 0x2A)
OFFWHITE     = RGBColor(0xED, 0xED, 0xED)
MUTED        = RGBColor(0x88, 0x88, 0x88)
ACCENT       = RGBColor(0xFF, 0x13, 0x00)
TIER1_GOLD   = RGBColor(0xFB, 0xBF, 0x24)
DETECT_GREEN = RGBColor(0x4A, 0xDE, 0x80)

FONT_DISPLAY = "Inter"     # Benzin substitute (system-safe)
FONT_MONO    = "JetBrains Mono"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL = 12


def set_bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=OFFWHITE, align=PP_ALIGN.LEFT, font=FONT_DISPLAY, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return box


def add_rect(slide, left, top, width, height, color, *, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_outlined_rect(slide, left, top, width, height, line_color, line_pt=2):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.background()
    shape.line.color.rgb = line_color
    shape.line.width = Pt(line_pt)
    shape.line.dash_style = 7  # MSO_LINE_DASH_STYLE.DASH
    return shape


def add_chrome(slide, slide_num):
    """Header + footer chrome shared across slides."""
    add_text(slide, Inches(0.5), Inches(0.3), Inches(4), Inches(0.4),
             "Spot the Brand.", size=14, bold=True, color=OFFWHITE)
    # red period dot
    add_rect(slide, Inches(2.06), Inches(0.55), Inches(0.08), Inches(0.08), ACCENT, rounded=True)
    add_text(slide, Inches(12.0), Inches(7.05), Inches(1.0), Inches(0.3),
             f"{slide_num:02d} / {TOTAL:02d}", size=10, color=MUTED,
             align=PP_ALIGN.RIGHT, font=FONT_MONO)
    add_text(slide, Inches(0.5), Inches(7.05), Inches(4), Inches(0.3),
             "spotyourbrand.com", size=10, color=MUTED, font=FONT_MONO)


def add_crosshair(slide):
    add_rect(slide, Inches(0), Inches(3.75), SLIDE_W, Pt(0.5), BORDER)
    add_rect(slide, Inches(6.66), Inches(0), Pt(0.5), SLIDE_H, BORDER)


def section_label(slide, text):
    add_text(slide, Inches(0.6), Inches(0.95), Inches(4), Inches(0.4),
             text, size=12, color=ACCENT, font=FONT_MONO, bold=True)


# ════════════════════════════════════════════════════════
# DECK v0.2 — BRAND LED
# ════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ─── 01 · Title ───
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    add_crosshair(s)
    add_outlined_rect(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(2.3), ACCENT)
    add_text(s, Inches(1.0), Inches(2.35), Inches(4), Inches(0.3),
             "DETECTED · 1.00", size=11, color=ACCENT, font=FONT_MONO, bold=True)
    add_text(s, Inches(1.2), Inches(2.9), Inches(11), Inches(1.4),
             "Spot the Brand.", size=88, bold=True, color=OFFWHITE)
    add_text(s, Inches(1.2), Inches(4.1), Inches(11), Inches(0.8),
             "Visual brand monitoring for the way people actually post.",
             size=24, color=OFFWHITE)
    add_text(s, Inches(1.2), Inches(5.4), Inches(11), Inches(0.5),
             "We see what your social tool misses.",
             size=18, color=MUTED, italic=True)
    # chrome (custom for title slide)
    add_text(s, Inches(0.5), Inches(0.3), Inches(4), Inches(0.4),
             "Spot the Brand.", size=14, bold=True, color=OFFWHITE)
    add_text(s, Inches(12.0), Inches(7.05), Inches(1.0), Inches(0.3),
             f"01 / {TOTAL:02d}", size=10, color=MUTED, align=PP_ALIGN.RIGHT, font=FONT_MONO)
    add_text(s, Inches(0.5), Inches(7.05), Inches(4), Inches(0.3),
             "by JackandAI · spotyourbrand.com", size=10, color=MUTED, font=FONT_MONO)

    # ─── 02 · Brand belief ───
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    add_chrome(s, 2)
    section_label(s, "WHAT WE BELIEVE")
    add_text(s, Inches(0.6), Inches(1.6), Inches(12), Inches(2.6),
             "Brands live in the wild.\nNot in hashtags.",
             size=64, bold=True, color=OFFWHITE)
    add_text(s, Inches(0.6), Inches(4.6), Inches(12), Inches(0.6),
             "In someone's hand at a festival. On a shelf behind a kitchen video.",
             size=20, color=OFFWHITE)
    add_text(s, Inches(0.6), Inches(5.2), Inches(12), Inches(0.6),
             "In a get-ready-with-me where the product is the punchline but never named.",
             size=20, color=OFFWHITE)
    add_text(s, Inches(0.6), Inches(5.9), Inches(12), Inches(0.6),
             "In a story that disappears in 24 hours but reaches 40.000 people.",
             size=20, color=OFFWHITE)
    add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
             "That's where brands actually live. Brand monitoring has been blind to it.",
             size=15, color=MUTED, italic=True)

    # ─── 03 · The gap ───
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    add_chrome(s, 3)
    section_label(s, "THE GAP")
    add_text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(1.0),
             "We see images. They read text.",
             size=52, bold=True, color=OFFWHITE)

    # Two big columns
    add_rect(s, Inches(0.6), Inches(3.2), Inches(6.0), Inches(3.4), SURFACE)
    add_rect(s, Inches(0.6), Inches(3.2), Inches(0.06), Inches(3.4), MUTED)
    add_text(s, Inches(0.85), Inches(3.4), Inches(5.5), Inches(0.4),
             "LEGACY TOOLING", size=11, color=MUTED, bold=True, font=FONT_MONO)
    add_text(s, Inches(0.85), Inches(3.8), Inches(5.5), Inches(0.6),
             "Reads captions, hashtags, @-mentions.",
             size=20, color=OFFWHITE)
    add_text(s, Inches(0.85), Inches(4.6), Inches(5.5), Inches(0.6),
             "Built when people wrote about brands.",
             size=15, color=MUTED)
    add_text(s, Inches(0.85), Inches(5.6), Inches(5.5), Inches(1.0),
             "Catches the ~20% of social mentions that come with a tag. Misses the rest.",
             size=15, color=MUTED, italic=True)

    add_rect(s, Inches(6.85), Inches(3.2), Inches(6.0), Inches(3.4), SURFACE)
    add_rect(s, Inches(6.85), Inches(3.2), Inches(0.06), Inches(3.4), ACCENT)
    add_text(s, Inches(7.1), Inches(3.4), Inches(5.5), Inches(0.4),
             "SPOT THE BRAND", size=11, color=ACCENT, bold=True, font=FONT_MONO)
    add_text(s, Inches(7.1), Inches(3.8), Inches(5.5), Inches(0.6),
             "Sees the product in the actual image.",
             size=20, color=OFFWHITE)
    add_text(s, Inches(7.1), Inches(4.6), Inches(5.5), Inches(0.6),
             "Built for how people post in 2026.",
             size=15, color=MUTED)
    add_text(s, Inches(7.1), Inches(5.6), Inches(5.5), Inches(1.0),
             "Catches the ~80% visible-but-untagged. Plus the tagged 20%. The full reality.",
             size=15, color=OFFWHITE, italic=True)

    # ─── 04 · Why now ───
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    add_chrome(s, 4)
    section_label(s, "WHY NOW")
    add_text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(1.6),
             "Computer vision\nfinally works for this.",
             size=52, bold=True, color=OFFWHITE)
    points = [
        ("Gemini Flash (2024-2025)", "5× cheaper than 2024 GPT-4V, 2× faster, equal accuracy on packshots."),
        ("Apify + adjacent infra", "Legitimate, scalable social scraping became commodity infra."),
        ("The category gap", "Brand monitoring stayed text-only for a decade. Window for vision-native is open."),
    ]
    for i, (head, body) in enumerate(points):
        top = Inches(4.0 + i * 0.95)
        add_rect(s, Inches(0.6), top + Inches(0.18), Inches(0.18), Inches(0.18), ACCENT, rounded=True)
        add_text(s, Inches(1.0), top, Inches(12), Inches(0.5),
                 head, size=18, bold=True, color=OFFWHITE)
        add_text(s, Inches(1.0), top + Inches(0.45), Inches(12), Inches(0.5),
                 body, size=14, color=MUTED)

    # ─── 05 · What we do ───
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    add_chrome(s, 5)
    section_label(s, "WHAT WE DO")
    add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.9),
             "Four steps. Automated. Daily.",
             size=30, bold=True, color=OFFWHITE)
    steps = [
        ("01", "Scrape",   "Daily IG + TikTok harvest by hashtag + creator. ~3k posts/day per brand."),
        ("02", "Detect",   "Computer vision finds your product packaging in the image. Bounding box per match."),
        ("03", "Verify",   "Higher-resolution second pass re-checks borderline cases. <5% false positive rate."),
        ("04", "Surface",  "Tier creators, score relevance, send Slack/email alerts, weekly PDF."),
    ]
    for i, (n, title, body) in enumerate(steps):
        top = Inches(2.6 + i * 1.0)
        add_text(s, Inches(0.6), top, Inches(0.8), Inches(0.7),
                 n, size=28, bold=True, color=ACCENT, font=FONT_MONO)
        add_text(s, Inches(1.5), top, Inches(2.2), Inches(0.7),
                 title, size=22, bold=True, color=OFFWHITE)
        add_text(s, Inches(4.0), top + Inches(0.05), Inches(9), Inches(0.7),
                 body, size=15, color=MUTED)

    # ─── 06 · Proof (customer example, framed carefully) ───
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    add_chrome(s, 6)
    section_label(s, "PROOF")
    add_text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(1.0),
             "Pilot brand. 90 days. Measured.",
             size=42, bold=True, color=OFFWHITE)
    add_text(s, Inches(0.6), Inches(2.5), Inches(12), Inches(0.5),
             "A Dutch hard seltzer brand has run Spot the Brand in production for a quarter. Numbers:",
             size=15, color=MUTED, italic=True)

    stats = [
        ("436", "verified product hits"),
        ("545", "creators identified"),
        ("86%", "had no #hashtag"),
        ("0.95", "avg confidence"),
    ]
    for i, (num, lbl) in enumerate(stats):
        left = Inches(0.6 + i * 3.05)
        add_rect(s, left, Inches(3.4), Inches(2.9), Inches(2.0), SURFACE)
        add_rect(s, left, Inches(3.4), Inches(0.06), Inches(2.0), ACCENT)
        add_text(s, left + Inches(0.3), Inches(3.65), Inches(2.5), Inches(1.0),
                 num, size=58, bold=True, color=OFFWHITE)
        add_text(s, left + Inches(0.3), Inches(4.85), Inches(2.5), Inches(0.5),
                 lbl.upper(), size=10, color=MUTED, bold=True, font=FONT_MONO)

    add_text(s, Inches(0.6), Inches(5.8), Inches(12), Inches(0.6),
             "Their old tool tracked 60 of the 436. We tracked the full reality.",
             size=18, color=OFFWHITE)
    add_text(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.5),
             "Your numbers will look different. The pattern — most mentions visible-but-untagged — holds across categories.",
             size=13, color=MUTED, italic=True)

    # ─── 07 · The product ───
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    add_chrome(s, 7)
    section_label(s, "THE PRODUCT")
    add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.9),
             "A multi-tenant SaaS. Built on Supabase + Vercel + Railway.",
             size=24, bold=True, color=OFFWHITE)
    features = [
        ("Live dashboard",         "Real-time feed of every detection with image, creator, confidence, tier."),
        ("Creator profiles",       "AI-scored relevance 0-10. Tier auto-promotion based on hit frequency."),
        ("Brand-specific refs",    "Upload your product packshots. We tune detection per brand."),
        ("Team + roles",           "Magic-link invites. Owner / admin / editor / viewer."),
        ("Alerts",                 "Slack webhook + email when a tier-1 creator lands new content."),
        ("Weekly PDF report",      "Designed insight report, emailed and downloadable."),
    ]
    for i, (k, v) in enumerate(features):
        col = i % 2
        row = i // 2
        left = Inches(0.6 + col * 6.2)
        top = Inches(2.6 + row * 1.45)
        add_rect(s, left, top, Inches(6), Inches(1.25), SURFACE)
        add_rect(s, left, top, Inches(0.06), Inches(1.25), ACCENT)
        add_text(s, left + Inches(0.3), top + Inches(0.15), Inches(5.6), Inches(0.4),
                 k, size=17, bold=True, color=OFFWHITE)
        add_text(s, left + Inches(0.3), top + Inches(0.6), Inches(5.6), Inches(0.6),
                 v, size=12, color=MUTED)

    # ─── 08 · Pricing ───
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    add_chrome(s, 8)
    section_label(s, "PRICING")
    add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.9),
             "Three tiers. Self-serve. 14-day free trial.",
             size=28, bold=True, color=OFFWHITE)

    plans = [
        ("STARTER",    "€500",   "/month",
         ["1 brand", "100 credits/mo", "Weekly report", "Monthly scans"]),
        ("PRO",        "€1.500", "/month",
         ["1 brand", "1.000 credits/mo", "Daily scans", "Pro verify", "Slack alerts", "Team seats"]),
        ("ENTERPRISE", "Custom", "",
         ["Multi-brand", "10k+ credits", "API access", "SLA", "Dedicated CSM", "White-label"]),
    ]
    for i, (name, price, suffix, features) in enumerate(plans):
        left = Inches(0.8 + i * 4.05)
        is_featured = name == "PRO"
        bg = ACCENT if is_featured else SURFACE
        text_col = OFFWHITE
        add_rect(s, left, Inches(2.6), Inches(3.85), Inches(4.2), bg)
        add_text(s, left + Inches(0.3), Inches(2.85), Inches(3.5), Inches(0.4),
                 name, size=13, bold=True, color=text_col, font=FONT_MONO)
        add_text(s, left + Inches(0.3), Inches(3.3), Inches(3.5), Inches(0.9),
                 price, size=46, bold=True, color=text_col)
        if suffix:
            add_text(s, left + Inches(0.3), Inches(4.25), Inches(3.5), Inches(0.4),
                     suffix, size=14, color=MUTED if not is_featured else OFFWHITE)
        for j, f in enumerate(features):
            top = Inches(4.7 + j * 0.32)
            add_text(s, left + Inches(0.3), top, Inches(3.4), Inches(0.3),
                     "· " + f, size=12, color=text_col)

    # ─── 09 · Competition ───
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    add_chrome(s, 9)
    section_label(s, "COMPETITION")
    add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.9),
             "We're the first vision-native player. Adjacent tools coexist.",
             size=26, bold=True, color=OFFWHITE)

    rows = [
        ("",                "Hashtag", "Visual", "Pricing",  "Backfill"),
        ("Storyclash",      "Yes",     "No",     "€1k+/mo",  "Limited"),
        ("Brand24 / Mention","Yes",    "No",     "€500+/mo", "No"),
        ("Tagger / Modash", "Yes",     "No",     "€2k+/mo",  "Manual"),
        ("Manual scroll",   "Maybe",   "Maybe",  "FTE cost", "No"),
        ("Spot the Brand",  "Yes",     "Yes",    "€500/mo",  "365 days"),
    ]
    col_w = [Inches(2.8), Inches(1.8), Inches(1.8), Inches(2.0), Inches(2.0)]
    col_x = [Inches(0.6 + sum(c.inches for c in col_w[:i])) for i in range(len(col_w))]
    for r_idx, row in enumerate(rows):
        top = Inches(2.7 + r_idx * 0.55)
        is_us = row[0] == "Spot the Brand"
        is_header = r_idx == 0
        bg = ACCENT if is_us else (SURFACE if r_idx % 2 == 0 else BLACK)
        if is_us or is_header or r_idx % 2 == 0:
            add_rect(s, Inches(0.6), top, Inches(10.4), Inches(0.55), bg)
        for c_idx, cell in enumerate(row):
            color = OFFWHITE if is_us or is_header else MUTED
            bold = is_us or is_header or c_idx == 0
            add_text(s, col_x[c_idx] + Inches(0.15), top + Inches(0.1), col_w[c_idx], Inches(0.4),
                     cell, size=12 if not is_header else 11, bold=bold, color=color,
                     font=FONT_MONO if is_header else FONT_DISPLAY)

    # ─── 10 · ICP ───
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    add_chrome(s, 10)
    section_label(s, "WHO BUYS THIS")
    add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.9),
             "Challenger brands with a visually distinctive product.",
             size=24, bold=True, color=OFFWHITE)

    icps = [
        ("Hard seltzer / RTD",     "Klar, Captain Morgan Spritz, hard-seltzer challengers entering NL"),
        ("Craft beer",             "Vandestreek, Lowlander, Brouwerij 't IJ, Charlie's Beer"),
        ("Functional drinks",      "RAW Superdrink, Crisp, Mother"),
        ("DTC personal care",      "Rituals (specific lines), about-time, Mother Mary, Beebop"),
        ("Specialty food",         "Tony's Chocolonely (per line), Stroopwafels by Lotte"),
        ("Spirits / cocktails",    "Hyke Gin, Klar, RTD cocktail brands"),
    ]
    for i, (cat, ex) in enumerate(icps):
        col = i % 2
        row = i // 2
        left = Inches(0.6 + col * 6.2)
        top = Inches(2.7 + row * 1.3)
        add_rect(s, left, top, Inches(6), Inches(1.15), SURFACE)
        add_rect(s, left, top, Inches(0.06), Inches(1.15), ACCENT)
        add_text(s, left + Inches(0.3), top + Inches(0.15), Inches(5.6), Inches(0.4),
                 cat, size=17, bold=True, color=OFFWHITE)
        add_text(s, left + Inches(0.3), top + Inches(0.55), Inches(5.6), Inches(0.5),
                 ex, size=12, color=MUTED)

    # ─── 11 · Traction ───
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    add_chrome(s, 11)
    section_label(s, "TRACTION")
    add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.9),
             "Built. Live. First customer in production.",
             size=28, bold=True, color=OFFWHITE)

    milestones = [
        ("Mar 2026",  "Pipeline v1 built",                    "Scraping + Gemini Flash + reference image fitting"),
        ("Apr 2026",  "Pilot brand in production",             "Dutch hard seltzer pilot. 10k credits, daily scans."),
        ("May 2026",  "SaaS infra: auth, billing, RLS",       "Supabase + Stripe + Vercel. Multi-tenant ready."),
        ("May 2026",  "Launch readiness reached",             "Demo public, Stripe wired, Customer Portal live."),
        ("Jun 2026",  "Target: 10 paying brands",             "€15k MRR. Break-even on infra cost."),
        ("Q4 2026",   "Target: 50 brands, multi-region",      "€75k MRR. Profitable."),
    ]
    for i, (date, t, sub) in enumerate(milestones):
        top = Inches(2.7 + i * 0.7)
        col = ACCENT if "Target" in t else OFFWHITE
        add_text(s, Inches(0.6), top, Inches(2), Inches(0.4),
                 date, size=12, color=MUTED, font=FONT_MONO, bold=True)
        add_text(s, Inches(2.6), top, Inches(5.5), Inches(0.4),
                 t, size=15, color=col, bold=True)
        add_text(s, Inches(8.0), top, Inches(5), Inches(0.4),
                 sub, size=13, color=MUTED, italic=True)

    # ─── 12 · Close ───
    s = prs.slides.add_slide(blank); set_bg(s, BLACK)
    add_chrome(s, 12)
    add_crosshair(s)
    add_outlined_rect(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.6), ACCENT)
    add_text(s, Inches(1.0), Inches(2.15), Inches(4), Inches(0.3),
             "DETECTED · 1.00", size=11, color=ACCENT, font=FONT_MONO, bold=True)
    add_text(s, Inches(1.2), Inches(2.6), Inches(11), Inches(1.4),
             "Let's see your brand.", size=72, bold=True, color=OFFWHITE)
    add_text(s, Inches(1.2), Inches(3.9), Inches(11), Inches(0.8),
             "14-day trial. No credit card. First hits within 24h.",
             size=22, color=OFFWHITE)
    add_text(s, Inches(1.2), Inches(5.5), Inches(11), Inches(0.5),
             "spotyourbrand.com", size=28, bold=True, color=ACCENT, font=FONT_MONO)
    add_text(s, Inches(1.2), Inches(6.1), Inches(11), Inches(0.5),
             "meinte@jackandai.com · jackandai.com",
             size=14, color=MUTED, font=FONT_MONO)

    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Wrote {OUTPUT_PATH}")


if __name__ == "__main__":
    build()
